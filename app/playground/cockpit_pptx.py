"""Cockpit da Conversa → apresentação PPTX (insumo para o Business Plan).

Recebe o SNAPSHOT que a UI computou (índices Q/R/C, gate, as três lentes —
Qualidade, Confiabilidade, Custo & TCO — e a conversa da sessão) e monta um deck
branded, slide a slide, com python-pptx. Sem template externo.

Princípio (igual ao cockpit): nenhum número sem dizer COMO foi obtido — cada slide
carrega os selos de proveniência e as ressalvas materiais (custo do juiz estimado,
self-hosted conta 0). É função PURA (dict → bytes) para ser testável isolada.

Os valores chegam CRUS (índices int|None; taxas 0–1|None; dinheiro em BRL float;
tokens/latência números) e são formatados aqui — o servidor não confia em strings
de exibição do cliente.
"""
from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# ── paleta (navy + azul de marca + semânticos) ──
INK = RGBColor(0x0E, 0x16, 0x26)
MUTED = RGBColor(0x69, 0x73, 0x88)
ACCENT = RGBColor(0x1D, 0x63, 0xD8)
GOOD = RGBColor(0x0F, 0x7A, 0x45)
WARN = RGBColor(0xC0, 0x7D, 0x00)
CRIT = RGBColor(0xD8, 0x32, 0x4C)
LINE = RGBColor(0xE2, 0xE6, 0xEE)
PANEL = RGBColor(0xF7, 0xF8, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9
_W = Emu(12192000)
_H = Emu(6858000)


def _f(x: Any) -> float | None:
    try:
        if x is None:
            return None
        v = float(x)
        return None if v != v else v  # NaN
    except (TypeError, ValueError):
        return None


def _brl(x: Any) -> str:
    v = _f(x)
    if v is None:
        return "—"
    s = f"{v:,.2f}"                    # 1,234.56 (en) → 1.234,56 (pt-BR)
    s = s.replace(",", "§").replace(".", ",").replace("§", ".")
    return f"R$ {s}"


def _brl_k(x: Any) -> str:
    v = _f(x)
    if v is None:
        return "—"
    if abs(v) >= 1000:
        s = f"{v / 1000:,.1f}".replace(",", "§").replace(".", ",").replace("§", ".")
        return f"R$ {s}k"
    return _brl(v)


def _usd(x: Any) -> str:
    v = _f(x)
    return "—" if v is None else f"US$ {v:.4f}"


def _pct(x: Any) -> str:
    v = _f(x)
    return "—" if v is None else f"{round(v * 100)}%"


def _score5(x: Any) -> str:
    v = _f(x)
    return "—" if v is None else f"{v * 5:.1f} /5"


def _idx(x: Any) -> str:
    v = _f(x)
    return "—" if v is None else f"{int(round(v))}"


def _int(x: Any) -> str:
    v = _f(x)
    if v is None:
        return "—"
    return f"{int(v):,}".replace(",", ".")


def _ms(x: Any) -> str:
    v = _f(x)
    if v is None:
        return "—"
    return f"{v / 1000:.1f} s" if v >= 1000 else f"{int(v)} ms"


def _txt(slide, l, t, w, h, text, *, size=14, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, spacing=1.0):
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def _rect(slide, l, t, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _header(slide, kicker, title):
    _txt(slide, 640080, 400000, 10911840, 300000, kicker, size=12, bold=True,
         color=ACCENT)
    _txt(slide, 640080, 640000, 10911840, 700000, title, size=30, bold=True)
    _rect(slide, 640080, 1360000, 10911840, 22000, LINE)


def _gate_color(k: str) -> RGBColor:
    return {"ok": GOOD, "warn": WARN, "crit": CRIT}.get(k or "", MUTED)


def _kv_lines(slide, l, t, w, rows, *, size=13, gap=430000):
    """rows: list of (label, value, value_color)."""
    y = t
    for label, value, vcolor in rows:
        _txt(slide, l, y, int(w * 0.62), gap, label, size=size, color=MUTED)
        _txt(slide, l + int(w * 0.62), y, int(w * 0.38), gap, value, size=size,
             bold=True, color=vcolor, align=PP_ALIGN.RIGHT)
        y += gap
    return y


def _prov_note(slide, text):
    _txt(slide, 640080, 6280000, 10911840, 360000, text, size=9, italic=True,
         color=MUTED)


def build_cockpit_pptx(snapshot: dict[str, Any]) -> bytes:
    """Monta o deck do Cockpit da Conversa a partir do snapshot da UI. Defensivo:
    chaves ausentes viram “—”, nunca quebram."""
    s = snapshot or {}
    idx = s.get("indices") or {}
    gate = s.get("gate") or {}
    q = s.get("quality") or {}
    ragas = q.get("ragas") or {}
    rel = s.get("reliability") or {}
    cost = s.get("cost") or {}
    tco = s.get("tco") or {}
    conv = s.get("conversation") or []
    pipeline = str(s.get("pipeline_name") or "—")
    session = str(s.get("session_id") or "")
    scope = "sessão" if (s.get("scope") == "sessao") else "este turno"
    generated = str(s.get("generated_at") or "")

    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H

    # ── Slide 1 — capa ──
    sl = _slide(prs)
    _rect(sl, 0, 0, int(_W), int(_H), INK)
    _txt(sl, 900000, 2200000, 10400000, 400000, "COCKPIT DA CONVERSA", size=15,
         bold=True, color=RGBColor(0x8F, 0xB4, 0xFF))
    _txt(sl, 900000, 2680000, 10400000, 1000000, pipeline, size=40, bold=True,
         color=WHITE)
    _txt(sl, 900000, 3720000, 10400000, 500000,
         "Avaliação big-picture — qualidade, confiabilidade e custo à luz do TCO",
         size=16, color=RGBColor(0xC0, 0xC9, 0xD8))
    meta = f"Recorte: {scope}"
    if session:
        meta += f"   ·   sessão {session[:8]}…"
    if generated:
        meta += f"   ·   {generated}"
    _txt(sl, 900000, 4360000, 10400000, 400000, meta, size=12,
         color=RGBColor(0x93, 0xA0, 0xB5))

    # ── Slide 2 — resumo executivo (gate + índices + TCO) ──
    sl = _slide(prs)
    _header(sl, "Resumo executivo", "Veredito, índices e TCO")
    gk = str(gate.get("k") or "")
    _rect(sl, 640080, 1560000, 10911840, 900000, PANEL, line=_gate_color(gk))
    _txt(sl, 820000, 1660000, 900000, 700000,
         {"ok": "✓", "warn": "!", "crit": "✕"}.get(gk, "•"), size=40, bold=True,
         color=_gate_color(gk), align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _txt(sl, 1720000, 1660000, 9700000, 360000, str(gate.get("t") or "—"),
         size=20, bold=True)
    _txt(sl, 1720000, 2020000, 9700000, 380000, str(gate.get("w") or ""),
         size=12, color=MUTED)

    trio = [("Qualidade", idx.get("q")), ("Confiabilidade", idx.get("r")),
            ("Eficiência", idx.get("c"))]
    cw, cgap = 3400000, 340000
    x0 = 640080
    for i, (label, val) in enumerate(trio):
        x = x0 + i * (cw + cgap)
        _rect(sl, x, 2680000, cw, 1500000, PANEL, line=LINE)
        vf = _f(val)
        col = MUTED if vf is None else (GOOD if vf >= 85 else WARN if vf >= 70 else CRIT)
        _txt(sl, x, 2820000, cw, 800000, _idx(val), size=48, bold=True, color=col,
             align=PP_ALIGN.CENTER)
        _txt(sl, x, 3760000, cw, 340000, label, size=14, bold=True,
             align=PP_ALIGN.CENTER, color=INK)

    band = tco.get("band_month") or {}
    yband = tco.get("band_year") or {}
    _txt(sl, 640080, 4460000, 6000000, 360000, "TCO mensal (banda)", size=13,
         color=MUTED)
    _txt(sl, 640080, 4760000, 6000000, 500000,
         f"{_brl_k(band.get('lo'))} – {_brl_k(band.get('hi'))}", size=26, bold=True)
    _txt(sl, 6800000, 4460000, 4700000, 360000, "Projeção anual", size=13,
         color=MUTED)
    _txt(sl, 6800000, 4760000, 4700000, 500000,
         f"{_brl_k(yband.get('lo'))} – {_brl_k(yband.get('hi'))}", size=26,
         bold=True, color=MUTED)
    _prov_note(sl, "Índices 0–100 · gate derivado dos índices · TCO projetado (banda −15%/+20%). "
                   "Selos e regras detalhados nos próximos slides.")

    # ── Slide 3 — Qualidade ──
    sl = _slide(prs)
    _header(sl, "① Qualidade da resposta · medido + heurístico",
            "Dimensões do juiz e RAGAS")
    if q.get("has_judge"):
        rows = [
            ("Factualidade (juiz)", _score5(q.get("factuality")), INK),
            ("Completude (juiz)", _score5(q.get("completeness")), INK),
            ("Aderência de tom (juiz)", _score5(q.get("tone")), INK),
            ("Segurança (juiz)", _score5(q.get("safety")), INK),
        ]
        _kv_lines(sl, 640080, 1700000, 5300000, rows, size=15, gap=560000)
        rrows = [
            ("faithfulness", _num2(ragas.get("faithfulness")), INK),
            ("answer relevancy", _num2(ragas.get("answer_relevancy")), INK),
            ("context precision", _num2(ragas.get("context_precision")), INK),
            ("context relevancy", _num2(ragas.get("context_relevancy")), INK),
        ]
        _rect(sl, 6300000, 1640000, 5251840, 2540000, PANEL, line=LINE)
        _txt(sl, 6500000, 1740000, 4900000, 340000, "RAGAS (heurístico)", size=13,
             bold=True, color=WARN)
        _kv_lines(sl, 6500000, 2160000, 4851840, rrows, size=13, gap=430000)
    else:
        _rect(sl, 640080, 1700000, 10911840, 700000, RGBColor(0xFB, 0xF0, 0xD8),
              line=WARN)
        _txt(sl, 820000, 1820000, 10500000, 500000,
             "O juiz (Verifier v2) não avaliou este recorte — Qualidade não medida. "
             "As demais lentes seguem válidas.", size=14, color=WARN,
             anchor=MSO_ANCHOR.MIDDLE)
    _txt(sl, 640080, 4360000, 10911840, 600000,
         "Não calculados (exigem casos-gold): groundedness · context recall · answer correctness.",
         size=12, italic=True, color=MUTED)
    _prov_note(sl, "faithfulness = factualidade/5; answer relevancy = completude/5; "
                   "context precision/relevancy = média do evidence_score. Proxy client-side, não o framework RAGAS.")

    # ── Slide 4 — Confiabilidade ──
    sl = _slide(prs)
    _header(sl, "② Confiabilidade & risco", "Por que confiar (ou não) na resposta")
    hall = rel.get("hallucinations")
    hall_v = _f(hall)
    rows = [
        ("Decisão final (FSM)", str(rel.get("fsm") or "—"), INK),
        ("Conformidade de contrato", _pct(rel.get("contract")),
         GOOD if _f(rel.get("contract")) == 1 else WARN if _f(rel.get("contract")) is not None else MUTED),
        ("Alegações sem lastro (alucinações)", _int(hall) if hall_v is not None else "—",
         CRIT if (hall_v or 0) > 0 else GOOD),
        ("Fundamentação (grounding)", _pct(rel.get("grounding")), INK),
        ("Cobertura do juiz (amostrado)", _pct(rel.get("coverage")), INK),
    ]
    _kv_lines(sl, 640080, 1740000, 10911840, rows, size=16, gap=740000)
    _prov_note(sl, "FSM/contrato/evidência: medido no trace. Cobertura amostrada — "
                   "cobertura < 100% significa que parte não foi auditada (não confundir com ausência de problema).")

    # ── Slide 5 — Custo & TCO ──
    sl = _slide(prs)
    _header(sl, "③ Custo & TCO · medido + projetado", "Unit economics e projeção")
    top = [
        ("Custo do recorte", _brl(cost.get("brl")), _usd(cost.get("usd"))),
        ("Tokens", _int(cost.get("tokens")), ""),
        ("Latência", _ms(cost.get("latency_ms")), ""),
    ]
    cw2 = 3500000
    for i, (label, big, sub) in enumerate(top):
        x = 640080 + i * (cw2 + 205920)
        _rect(sl, x, 1580000, cw2, 900000, PANEL, line=LINE)
        _txt(sl, x + 160000, 1650000, cw2 - 320000, 300000, label, size=11, color=MUTED)
        _txt(sl, x + 160000, 1930000, cw2 - 320000, 400000, big, size=20, bold=True)
        if sub:
            _txt(sl, x + 160000, 2300000, cw2 - 320000, 240000, sub, size=10, color=MUTED)
    # QAC
    _rect(sl, 640080, 2620000, 5300000, 780000, RGBColor(0xE7, 0xEF, 0xFC))
    _txt(sl, 820000, 2700000, 5000000, 400000, _brl(cost.get("qac_brl")), size=24,
         bold=True, color=ACCENT)
    _txt(sl, 820000, 3120000, 5000000, 260000, "por resposta CONFIÁVEL (custo ÷ aprovação)",
         size=11, color=ACCENT)
    # escada
    ladder = [
        ("Inferência LLM", tco.get("inf_month"), ACCENT),
        ("Juiz / verificador (estimado)", tco.get("judge_month"), WARN),
        ("Ferramentas SaaS (declarado)", tco.get("saas_month"), GOOD),
        ("Infra self-hosted (não medido)", None, MUTED),
    ]
    y = 2620000
    for label, val, col in ladder:
        _txt(sl, 6300000, y, 3600000, 320000, label, size=11, color=MUTED)
        _txt(sl, 9900000, y, 1651920, 320000,
             "não medido" if val is None else _brl_k(val), size=12, bold=True,
             color=col, align=PP_ALIGN.RIGHT)
        y += 360000
    _txt(sl, 6300000, y + 60000, 3600000, 320000, "TCO mensal — banda", size=12, bold=True)
    _txt(sl, 9200000, y + 60000, 2351920, 320000,
         f"{_brl_k(band.get('lo'))}–{_brl_k(band.get('hi'))}", size=13, bold=True,
         align=PP_ALIGN.RIGHT)
    prem = (f"Premissas: volume {_int(tco.get('volume'))} conv./mês · "
            f"câmbio {_brl(tco.get('fx'))} · meta {_brl(tco.get('alvo'))}/resposta.")
    _txt(sl, 640080, 4360000, 10911840, 340000, prem, size=12, color=INK)
    _prov_note(sl, "Ressalvas: preço LLM = tabela interna 2026-05; custo do juiz ESTIMADO (não instrumentado); "
                   "self-hosted conta 0 (GPU/infra fora do modelo); input pode subcontar em turnos multi-chamada.")

    # ── Slide 6+ — a conversa ──
    _conversation_slides(prs, conv)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _num2(x: Any) -> str:
    v = _f(x)
    return "—" if v is None else f"{v:.2f}"


def _conversation_slides(prs, conv: list) -> None:
    """A conversa da sessão — pagina em quantos slides forem necessários."""
    def new():
        sl = _slide(prs)
        _header(sl, "A conversa da sessão", "Transcrição")
        return sl, 1620000

    if not conv:
        sl, _ = new()
        _txt(sl, 640080, 1800000, 10911840, 400000, "(sem turnos)", size=14,
             color=MUTED)
        return

    sl, y = new()
    limit = 6200000
    for m in conv:
        role = "Você" if (m or {}).get("role") == "user" else "Pipeline"
        text = str((m or {}).get("text") or "")
        if len(text) > 900:
            text = text[:900] + "…"
        rc = ACCENT if role == "Você" else INK
        lines = 1 + max(0, len(text)) // 95
        h = 300000 + lines * 210000
        if y + h > limit:
            sl, y = new()
        _txt(sl, 640080, y, 1500000, 260000, role, size=11, bold=True, color=rc)
        _txt(sl, 2200000, y, 9351920, h, text, size=12, color=INK, spacing=1.05)
        y += h + 140000
