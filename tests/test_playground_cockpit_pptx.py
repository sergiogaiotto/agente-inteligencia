"""Export do Cockpit da Conversa em PPTX.

O deck é montado server-side a partir do SNAPSHOT que a UI computou (índices, lentes
e TCO já no recorte) + a conversa da sessão. Testamos o gerador PURO (dict → bytes,
sem I/O) e o cabeamento (rota + método/botão no template).
"""
import io
from pathlib import Path

import pytest

from app.playground.cockpit_pptx import build_cockpit_pptx

PG = Path("app/templates/pages/mesh_playground.html")
ROUTE = Path("app/routes/playground.py")

SNAP = {
    "pipeline_name": "Aurora — Crédito PJ",
    "session_id": "9049e2bd-aaaa-bbbb-cccc-ddddeeeeffff",
    "scope": "sessao",
    "generated_at": "11/07/2026 10:30",
    "indices": {"q": 82, "r": 75, "c": None},
    "gate": {"k": "warn", "t": "Apto com ressalvas",
             "w": "5 alegações sem lastro · cobertura do juiz 50%"},
    "quality": {"has_judge": True, "factuality": 0.60, "completeness": 1.0,
                "tone": 1.0, "safety": 1.0,
                "ragas": {"faithfulness": 0.60, "answer_relevancy": 1.0,
                          "context_precision": 0.93, "context_relevancy": 0.93}},
    "reliability": {"fsm": "LogAndClose", "contract": 1.0, "hallucinations": 5,
                    "grounding": 1.0, "coverage": 0.5},
    "cost": {"brl": 0.0, "usd": 0.0, "tokens": 8856, "latency_ms": 28200,
             "qac_brl": 0.0},
    "tco": {"volume": 1000, "fx": 5.3, "alvo": 0.5, "saas": 0,
            "inf_month": 0.0, "judge_month": 0.0, "saas_month": 0.0,
            "band_month": {"lo": 0.0, "hi": 0.0},
            "band_year": {"lo": 0.0, "hi": 0.0}},
    "conversation": [
        {"role": "user", "text": "Preciso de ajuda com meu limite de crédito."},
        {"role": "assistant", "text": "Claro — seu perfil comporta elevação [E1]."},
    ],
}


def _all_text(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    out = []
    for slide in prs.slides:
        for shp in slide.shapes:
            if shp.has_text_frame:
                out.append(shp.text_frame.text)
    return "\n".join(out), len(prs.slides)


def test_pptx_valido_multi_slide_e_conteudo():
    data = build_cockpit_pptx(SNAP)
    assert isinstance(data, bytes) and len(data) > 5000
    assert data[:2] == b"PK"                     # OOXML = zip

    text, n_slides = _all_text(data)
    assert n_slides >= 6                          # capa + resumo + 3 lentes + conversa
    # identidade + índices + lentes
    assert "COCKPIT DA CONVERSA" in text
    assert "Aurora — Crédito PJ" in text
    assert "82" in text and "75" in text          # índices Q/R
    assert "Qualidade" in text and "Confiabilidade" in text and "TCO" in text
    # honestidade: RAGAS heurístico + o que falta + ressalvas de custo
    assert "RAGAS" in text
    assert "groundedness" in text and "context recall" in text
    assert "self-hosted conta 0" in text
    assert "por resposta CONFIÁVEL" in text
    # a conversa da sessão foi incluída
    assert "limite de crédito" in text


def test_pptx_defensivo_com_snapshot_vazio():
    """Snapshot vazio não pode quebrar — tudo vira “—”, deck ainda sai."""
    data = build_cockpit_pptx({})
    assert data[:2] == b"PK"
    _text, n = _all_text(data)
    assert n >= 6


def test_pptx_formata_pt_br_e_recorte():
    data = build_cockpit_pptx(SNAP)
    text, _ = _all_text(data)
    assert "R$" in text                          # dinheiro em BRL
    assert "sessão" in text                       # recorte legível
    assert "9049e2bd" in text                     # sessão abreviada na capa


def test_rota_pptx_registrada():
    src = ROUTE.read_text(encoding="utf-8")
    assert '@router.post("/cockpit/pptx")' in src
    assert "build_cockpit_pptx" in src
    assert "presentationml.presentation" in src            # media type PPTX
    assert "Content-Disposition" in src and "attachment" in src
    assert "Depends(require_user)" in src                  # auth por usuário


def test_botao_e_metodo_no_template():
    src = PG.read_text(encoding="utf-8")
    assert 'data-testid="pg-cockpit-pptx"' in src
    assert "exportCockpitPptx()" in src and "_cockpitSnapshot()" in src
    # o snapshot leva os índices, as lentes E a conversa da sessão
    assert "indices:" in src and "quality:" in src and "reliability:" in src and "tco:" in src
    assert "conversation: this.chat.map(" in src
    # POST cookie (ação de UI) para a rota do playground
    assert "'/api/v1/playground/cockpit/pptx'" in src
    assert "a.download = 'cockpit-'" in src
